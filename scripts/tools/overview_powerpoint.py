"""
Generate a PowerPoint presentation of the most recent figures shown in the
IOOS Model Comparisons Archive Explorer 'Overview' tab.

One slide per region/variable/depth combination, with both model-comparison
images side-by-side where available (RTOFS vs CMEMS, RTOFS vs ESPC).
Guam and Fiji only produce the ESPC vs CMEMS image, mirroring the web app.

Usage:
    python overview_powerpoint.py
    python overview_powerpoint.py --output my_overview.pptx
    python overview_powerpoint.py --regions Caribbean "Gulf of Mexico"

Dependencies (install if missing):
    pip install python-pptx requests
"""
import argparse
import datetime
import io
import re
import sys
from pathlib import Path

try:
    import requests
except ImportError:
    sys.exit("requests is required:  pip install requests")

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    sys.exit("python-pptx is required:  pip install python-pptx")

# ── Metadata ───────────────────────────────────────────────────────────────

# Default region order for the presentation (products are discovered dynamically).
OVERVIEW_REGIONS = [
    "Caribbean",
    "Gulf of Mexico",
    "South Atlantic Bight",
    "Mid Atlantic Bight",
    "West Florida Shelf",
    "Tropical Western Atlantic",
    "Eastern Pacific - Mexico",
    "Hawaii",
    "Guam",
    "Fiji",
]

# Variable display order; depths within each variable are sorted numerically.
VARIABLE_ORDER = ["temperature", "salinity", "ocean_heat_content", "currents"]

# Regions that show only the ESPC vs CMEMS comparison (no RTOFS images)
ESPC_CMEMS_ONLY = {"Guam", "Fiji"}

REGION_FOLDER_MAP = {
    "Eastern Pacific - Mexico": "mexico_pacific",
}

BASE_URL      = "https://rucool.marine.rutgers.edu/hurricane/model_comparisons/maps"
BASE_PROFILES = "https://rucool.marine.rutgers.edu/hurricane/model_comparisons/profiles"

# Slide colours
NAVY  = RGBColor(0x0d, 0x2b, 0x5e)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY  = RGBColor(0x55, 0x55, 0x55)


# ── Server helpers ─────────────────────────────────────────────────────────

def _region_url_key(region: str) -> str:
    return REGION_FOLDER_MAP.get(region, region.replace(" ", "_").lower())


def _product_sort_key(product: str) -> tuple[int, int]:
    """Sort by VARIABLE_ORDER index, then by depth numerically ascending."""
    for i, var in enumerate(VARIABLE_ORDER):
        if product == var or product.startswith(var + "_"):
            m = re.search(r'(\d+)m$', product)
            return (i, int(m.group(1)) if m else 0)
    return (len(VARIABLE_ORDER), 0)


def _product_label(product: str) -> str:
    """Human-readable label: 'temperature_200m' → 'Temperature 200m'."""
    parts = product.split("_")
    if re.match(r'^\d+m$', parts[-1]):
        # Keep depth token lowercase (e.g. '200m', not '200M')
        return " ".join(p.title() for p in parts[:-1]) + " " + parts[-1]
    return " ".join(p.title() for p in parts)


def discover_products(region: str) -> list[str]:
    """
    Scrape the region's maps directory on the server and return every
    variable_depth folder that exists, sorted by variable then depth.

    This means the presentation always reflects what is actually published —
    no hardcoded depth lists required.
    """
    key  = _region_url_key(region)
    base = f"{BASE_URL}/{key}/"
    try:
        r = requests.get(base, timeout=10)
        if r.status_code != 200:
            return []
        folders = re.findall(r'href="([a-z][a-z_0-9]*)/"', r.text)
        products = [
            f for f in folders
            if any(f == v or f.startswith(v + "_") for v in VARIABLE_ORDER)
        ]
        return sorted(set(products), key=_product_sort_key)
    except Exception as exc:
        print(f"  Warning: could not discover products for {region}: {exc}")
        return []


def get_latest_map_info(region: str, variable_depth: str,
                        date_filter: str | None = None):
    """
    Scrape the server directory listing to find the most recent date and
    6-hourly time step for the given region / variable-depth product.

    If date_filter is "YYYY-MM-DD", returns the latest time available for
    that specific date rather than the overall most recent date.

    Returns (date_str "YYYY-MM-DD", time_str "HHZ") or None.
    """
    key  = _region_url_key(region)
    base = f"{BASE_URL}/{key}/{variable_depth}/"
    try:
        if date_filter:
            year  = date_filter[:4]
            month = date_filter[5:7]
            r = requests.get(f"{base}{year}/{month}/", timeout=10)
            matches = re.findall(r'(\d{4}-\d{2}-\d{2})T(\d{6})Z', r.text)
            matches = [(d, t) for d, t in matches if d == date_filter]
        else:
            r = requests.get(base, timeout=10)
            years = sorted(re.findall(r'href="(\d{4})/"', r.text), reverse=True)
            if not years:
                return None
            year = years[0]
            r = requests.get(f"{base}{year}/", timeout=10)
            months = sorted(re.findall(r'href="(\d{2})/"', r.text), reverse=True)
            if not months:
                return None
            month = months[0]
            r = requests.get(f"{base}{year}/{month}/", timeout=10)
            matches = re.findall(r'(\d{4}-\d{2}-\d{2})T(\d{6})Z', r.text)

        if not matches:
            return None
        date_str, time_code = sorted(set(matches), reverse=True)[0]
        return date_str, time_code[:2] + "Z"
    except Exception as exc:
        print(f"  Warning: could not scrape {region}/{variable_depth}: {exc}")
        return None


def build_map_urls(region: str, variable_depth: str,
                   date_obj: datetime.date, time_str: str):
    """
    Return (url_copernicus, url_espc, url_espc_cmems, url_goes).
    Mirrors build_map_urls() in webapps/app.py exactly.
    url_goes is None for OHC and ESPC-only regions.
    """
    vd_mod  = variable_depth.replace("_", "-")
    year    = date_obj.year
    month   = f"{date_obj.month:02d}"
    day     = f"{date_obj.day:02d}"
    hour    = time_str.replace("Z", "").zfill(2)
    tcode   = f"{hour}0000"
    key     = _region_url_key(region)
    slug    = key if variable_depth.startswith("currents") else key.replace("_", "-")
    folder  = f"{BASE_URL}/{key}/{variable_depth}/{year}/{month}/"

    if variable_depth == "ocean_heat_content":
        if region in ESPC_CMEMS_ONLY:
            fname = f"{region}_{year}-{month}-{day}T{tcode}Z_ohc_rtofs-espc-cmems.png"
            return None, None, f"{folder}{fname}", None
        else:
            stem = f"{key}_{year}-{month}-{day}T{tcode}Z"
            return (
                f"{folder}{stem}_heat_content_rtofs-cmems.png",
                f"{folder}{stem}_heat_content_rtofs-espc.png",
                f"{folder}{stem}_heat_content_espc-cmems.png",
                None,
            )
    else:
        stem = f"{slug}_{year}-{month}-{day}T{tcode}Z_{vd_mod}"
        return (
            f"{folder}{stem}_rtofs-vs-cmems.png",
            f"{folder}{stem}_rtofs-vs-espc.png",
            f"{folder}{stem}_espc-vs-cmems.png",
            f"{folder}{stem}_rtofs-vs-GOES.png",
        )


def fetch_image(url: str | None) -> io.BytesIO | None:
    """Download an image and return a BytesIO buffer, or None if unavailable."""
    if not url:
        return None
    try:
        r = requests.get(url, timeout=20)
        ct = r.headers.get("Content-Type", "")
        if r.status_code == 200 and ct.startswith("image"):
            return io.BytesIO(r.content)
    except Exception:
        pass
    return None


# ── Slide builders ─────────────────────────────────────────────────────────

def _text_box(slide, left, top, width, height, text, font_size, bold=False,
              color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tb


def add_title_slide(prs: Presentation, subtitle: str):
    layout = prs.slide_layouts[6]  # blank
    slide  = prs.slides.add_slide(layout)
    W, H   = prs.slide_width, prs.slide_height

    # Navy background
    bg = slide.shapes.add_shape(
        1, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    _text_box(slide, Inches(0.5), Inches(2.8), W - Inches(1), Inches(1.2),
              "IOOS Model Comparisons", 40, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER)
    _text_box(slide, Inches(0.5), Inches(4.1), W - Inches(1), Inches(0.7),
              "Overview — Most Recent Figures", 24, color=RGBColor(0xAA, 0xCC, 0xFF),
              align=PP_ALIGN.CENTER)
    _text_box(slide, Inches(0.5), Inches(5.0), W - Inches(1), Inches(0.5),
              subtitle, 16, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)


def add_section_slide(prs: Presentation, region: str):
    layout = prs.slide_layouts[6]  # blank
    slide  = prs.slides.add_slide(layout)
    W, H   = prs.slide_width, prs.slide_height

    # Dark accent bar on left
    bar = slide.shapes.add_shape(1, 0, 0, Inches(0.4), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    _text_box(slide, Inches(0.7), Inches(3.5), W - Inches(1), Inches(1.2),
              region, 40, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    _text_box(slide, Inches(0.7), Inches(4.7), W - Inches(1), Inches(0.5),
              "Model Comparison Figures  ·  Most Recent Available", 14,
              color=GREY, align=PP_ALIGN.LEFT)


def add_image_slide(prs: Presentation, region: str, product_label: str,
                    images: list[tuple[io.BytesIO | None, str]],
                    date_str: str, time_str: str):
    """
    Add a content slide.

    images: list of (BytesIO_or_None, label_string)
    """
    layout = prs.slide_layouts[6]  # blank
    slide  = prs.slides.add_slide(layout)
    W, H   = prs.slide_width, prs.slide_height

    # Header bar
    hdr_h = Inches(0.55)
    hdr = slide.shapes.add_shape(1, 0, 0, W, hdr_h)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = NAVY
    hdr.line.fill.background()

    _text_box(slide, Inches(0.15), 0, W - Inches(3), hdr_h,
              f"{region}  ·  {product_label}", 13, bold=True, color=WHITE)
    _text_box(slide, W - Inches(2.8), 0, Inches(2.7), hdr_h,
              f"{date_str}  {time_str}", 12, color=RGBColor(0xAA, 0xCC, 0xFF),
              align=PP_ALIGN.RIGHT)

    # Image area
    margin   = Inches(0.12)
    lbl_h    = Inches(0.28)
    top      = hdr_h + margin
    avail_h  = H - top - margin
    avail_w  = W - 2 * margin

    valid = [(buf, lbl) for buf, lbl in images if buf is not None]
    n = len(valid)
    if n == 0:
        return

    img_w = (avail_w - (n - 1) * margin) / n
    img_h = avail_h - lbl_h

    for i, (buf, label) in enumerate(valid):
        left = margin + i * (img_w + margin)
        # Label
        lbl_box = slide.shapes.add_textbox(left, top, img_w, lbl_h)
        lbl_tf  = lbl_box.text_frame
        lbl_p   = lbl_tf.paragraphs[0]
        lbl_p.text = label
        lbl_p.alignment = PP_ALIGN.CENTER
        lbl_run = lbl_p.runs[0]
        lbl_run.font.size = Pt(10)
        lbl_run.font.bold = True
        lbl_run.font.color.rgb = GREY
        # Image — scale to fit slot without stretching, then centre
        try:
            pic = slide.shapes.add_picture(buf, 0, 0)  # natural size first
            scale  = min(img_w / pic.width, img_h / pic.height)
            fit_w  = int(pic.width  * scale)
            fit_h  = int(pic.height * scale)
            pic.left   = int(left + (img_w - fit_w) / 2)
            pic.top    = int(top + lbl_h + (img_h - fit_h) / 2)
            pic.width  = fit_w
            pic.height = fit_h
        except Exception as exc:
            print(f"  Warning: could not embed image '{label}': {exc}")


# ── Glider profile helpers ─────────────────────────────────────────────────

def get_latest_glider_date() -> datetime.date | None:
    """Scrape the glider profiles directory for the most recent dated folder."""
    base = f"{BASE_PROFILES}/gliders/"
    try:
        r = requests.get(base, timeout=10)
        years = sorted(re.findall(r'href="(\d{4})/"', r.text), reverse=True)
        if not years:
            return None
        year = years[0]
        r = requests.get(f"{base}{year}/", timeout=10)
        month_days = sorted(re.findall(r'href="(\d{2}-\d{2})/"', r.text), reverse=True)
        if not month_days:
            return None
        month_str, day_str = month_days[0].split("-")
        return datetime.date(int(year), int(month_str), int(day_str))
    except Exception as exc:
        print(f"  Warning: could not find latest glider date: {exc}")
        return None


def get_glider_files(date_obj: datetime.date) -> dict[str, str]:
    """
    Scrape the glider profiles directory for date_obj and return
    {glider_id: full_url} for every 400m PNG found.

    Using the actual filenames from the server listing avoids any
    mismatch between a reconstructed URL and what was actually saved.
    """
    year      = date_obj.year
    month_day = f"{date_obj.month:02d}-{date_obj.day:02d}"
    base_url  = f"{BASE_PROFILES}/gliders/{year}/{month_day}/"
    try:
        r = requests.get(base_url, timeout=8)
        if r.status_code != 200:
            return {}
        result: dict[str, str] = {}
        for m in re.finditer(r'href="([^"]+_400m\.png)"', r.text):
            fname = m.group(1)
            gid   = fname.split("_20")[0]
            result[gid] = f"{base_url}{fname}"
        return result
    except Exception as exc:
        print(f"  Warning: could not get glider files for {date_obj}: {exc}")
        return {}


# ── Main ───────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Build a PowerPoint of IOOS Model Comparisons overview figures."
    )
    parser.add_argument(
        "--output", "-o",
        default="ioos_model_comparisons_overview.pptx",
        help="Output .pptx filename (default: ioos_model_comparisons_overview.pptx)",
    )
    parser.add_argument(
        "--regions", nargs="+", metavar="REGION", default=None,
        help="Limit to specific region(s) (default: all overview regions)",
    )
    parser.add_argument(
        "--date", "-d", metavar="YYYY-MM-DD", default=None,
        help="Use figures from a specific date instead of the most recent available",
    )
    args = parser.parse_args()

    if args.date:
        try:
            datetime.datetime.strptime(args.date, "%Y-%m-%d")
        except ValueError:
            sys.exit(f"Invalid --date value '{args.date}' — expected YYYY-MM-DD")

    regions  = args.regions or OVERVIEW_REGIONS
    ref_date = (
        datetime.datetime.strptime(args.date, "%Y-%m-%d").date()
        if args.date
        else datetime.date.today()
    )

    prs = Presentation()
    prs.slide_width  = Inches(16)
    prs.slide_height = Inches(9)

    today = datetime.date.today().strftime("%Y-%m-%d")
    subtitle = f"Figures for {args.date}  ·  Generated {today}" if args.date else f"Generated {today}"
    add_title_slide(prs, subtitle)

    total_slides = 0

    for region in regions:
        print(f"\n{'='*60}\n  Region: {region}")

        products = discover_products(region)
        if not products:
            print("  No products found on server — skipping.")
            continue

        add_section_slide(prs, region)
        espc_only = region in ESPC_CMEMS_ONLY

        for product in products:
            label = _product_label(product)
            print(f"  {product:<35}", end="", flush=True)

            latest = get_latest_map_info(region, product, date_filter=args.date)
            if not latest:
                print("not found")
                continue

            date_str, time_str = latest
            date_obj = datetime.datetime.strptime(date_str, "%Y-%m-%d").date()
            if (ref_date - date_obj).days > 1:
                print(f"stale ({date_str})")
                continue
            url_cop, url_espc, url_espc_cmems, url_goes = build_map_urls(
                region, product, date_obj, time_str
            )

            if espc_only:
                buf = fetch_image(url_espc_cmems)
                pairs = [(buf, "ESPC vs CMEMS")]
            else:
                pairs = [
                    (fetch_image(url_cop),  "RTOFS vs CMEMS (Copernicus)"),
                    (fetch_image(url_espc), "RTOFS vs ESPC"),
                    (fetch_image(url_goes), "RTOFS vs GOES SST"),
                ]

            n_ok = sum(1 for buf, _ in pairs if buf is not None)
            if n_ok == 0:
                print("images unavailable")
                continue

            for buf, img_label in pairs:
                if buf is None:
                    continue
                add_image_slide(prs, region, f"{label}  ·  {img_label}", [(buf, img_label)], date_str, time_str)
                total_slides += 1
            print(f"ok  ({date_str} {time_str}, {n_ok} image(s))")

    # ── Glider profiles ───────────────────────────────────────────────────
    print(f"\n{'='*60}\n  Glider Profiles")
    if args.date:
        glider_date = datetime.datetime.strptime(args.date, "%Y-%m-%d").date()
    else:
        glider_date = get_latest_glider_date()

    if not glider_date:
        print("  Could not determine latest glider date")
    elif (ref_date - glider_date).days > 1:
        print(f"  Glider data is stale ({glider_date}, {(ref_date - glider_date).days} days old) — skipping")
    else:
        print(f"  Date: {glider_date}")
        glider_files = get_glider_files(glider_date)
        if glider_files:
            add_section_slide(prs, "Glider Profiles")
            date_str = glider_date.strftime("%Y-%m-%d")
            for gid, url in glider_files.items():
                buf = fetch_image(url)
                status = "ok" if buf else "unavailable"
                print(f"  {gid:<40} {status}")
                if buf:
                    add_image_slide(prs, "Glider Profiles", gid, [(buf, gid)], date_str, "")
                    total_slides += 1
        else:
            print(f"  No glider profiles found for {glider_date}")

    output = Path(args.output)
    prs.save(output)
    print(f"\n{'='*60}")
    print(f"Saved {total_slides} content slides → {output.resolve()}")


if __name__ == "__main__":
    main()
